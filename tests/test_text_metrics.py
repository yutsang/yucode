"""text_metrics: real glyph-width text fitting for PPTX shapes, ported from a
sibling project's fdd_utils/text_metrics.py. Exists because the fdd-commentary
skill's original chars-per-inch capacity heuristic was too rough to reliably
tell an LLM "this slot has room for 3 more sentences" -- a real run produced
commentary using only 1-2 sentences of slots sized for much more.
"""
from __future__ import annotations

import pytest

from coding_agent.tools import text_metrics


def _skip_if_no_font(is_cjk: bool) -> None:
    if text_metrics.resolve_font_path("Arial", is_cjk=is_cjk) is None:
        pytest.skip("no resolvable font on this machine")


class TestFontResolution:
    def test_resolves_a_latin_font_on_this_machine(self) -> None:
        _skip_if_no_font(False)
        path = text_metrics.resolve_font_path("Arial", is_cjk=False)
        assert path is not None

    def test_missing_font_raises_filenotfound_not_a_crash(self, monkeypatch) -> None:
        # resolve_font_path always tries the OS fallback chain regardless of
        # the family string, so a real font is found on any normal dev/CI/
        # Windows box even for a nonsense family name -- force the "nothing
        # resolves anywhere" branch directly instead. Clear _FONT_CACHE and use
        # a size no other test in this file uses, so a cache hit from another
        # test can't mask resolve_font_path never being consulted.
        text_metrics._FONT_CACHE.clear()
        monkeypatch.setattr(text_metrics, "resolve_font_path", lambda family, *, is_cjk: None)
        with pytest.raises(FileNotFoundError):
            text_metrics.get_font("Arial", 12.345, is_cjk=False)


class TestChineseDetection:
    def test_english_is_not_predominantly_chinese(self) -> None:
        assert text_metrics.is_predominantly_chinese("Cash increased to CNY7.9 million.") is False

    def test_english_naming_a_chinese_entity_is_still_not_predominantly_chinese(self) -> None:
        # A single CJK name inside mostly-Latin prose must not flip the whole
        # paragraph to CJK wrap rules -- this is exactly the false-positive
        # the >30% ratio (not "contains any CJK char") threshold guards against.
        text = "The related-party loan was owed to 维彧 and settled in full."
        assert text_metrics.is_predominantly_chinese(text) is False

    def test_chinese_paragraph_is_predominantly_chinese(self) -> None:
        assert text_metrics.is_predominantly_chinese("现金及银行存款由2024年12月31日增加至2025年") is True

    def test_empty_string_is_not_chinese(self) -> None:
        assert text_metrics.is_predominantly_chinese("") is False


class TestKinsoku:
    def test_pulls_leading_forbidden_punctuation_back(self) -> None:
        lines = ["第一行的內容", "。第二行"]
        result = text_metrics._apply_kinsoku(lines)
        assert result[0].endswith("。")
        assert not result[1].startswith("。")

    def test_single_line_is_unchanged(self) -> None:
        assert text_metrics._apply_kinsoku(["only one line"]) == ["only one line"]


class TestWrapParagraph:
    def test_wraps_long_latin_text_into_multiple_lines(self) -> None:
        _skip_if_no_font(False)
        font = text_metrics.get_font("Arial", 9.0, is_cjk=False)
        text = "the balance as at 31 January 2026 represented CNY7.9 million of cash at bank"
        lines = text_metrics.wrap_paragraph(text, font, max_width_pt=100.0)
        assert len(lines) > 1
        assert "".join(lines).replace(" ", "") == text.replace(" ", "")

    def test_short_text_fits_on_one_line_given_enough_width(self) -> None:
        _skip_if_no_font(False)
        font = text_metrics.get_font("Arial", 9.0, is_cjk=False)
        lines = text_metrics.wrap_paragraph("short text", font, max_width_pt=1000.0)
        assert lines == ["short text"]

    def test_blank_paragraph_returns_one_empty_line(self) -> None:
        _skip_if_no_font(False)
        font = text_metrics.get_font("Arial", 9.0, is_cjk=False)
        assert text_metrics.wrap_paragraph("", font, max_width_pt=100.0) == [""]


class TestEstimateCapacity:
    @pytest.fixture()
    def text_shape(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(1), Inches(1))
        return shape

    def test_short_text_in_a_large_box_has_low_fill_ratio(self, text_shape) -> None:
        _skip_if_no_font(False)
        result = text_metrics.estimate_capacity(
            text_shape, ["Cash increased."], is_chinese=False,
        )
        assert result["fits"] is True
        assert result["fill_ratio"] < 0.5
        assert result["remaining_lines_estimate"] > 0

    def test_long_text_in_a_small_box_overflows(self, text_shape) -> None:
        _skip_if_no_font(False)
        long_text = " ".join(["word"] * 300)
        result = text_metrics.estimate_capacity(text_shape, [long_text], is_chinese=False)
        assert result["fits"] is False
        assert result["fill_ratio"] > 1.0

    def test_auto_detects_chinese_when_is_chinese_omitted(self, text_shape) -> None:
        _skip_if_no_font(True)
        result = text_metrics.estimate_capacity(
            text_shape, ["现金及银行存款由2024年12月31日增加至2025年"],
        )
        assert result["is_chinese"] is True

    def test_multiple_paragraphs_each_add_a_gap(self, text_shape) -> None:
        _skip_if_no_font(False)
        one_para = text_metrics.estimate_capacity(text_shape, ["one line"], is_chinese=False)
        two_para = text_metrics.estimate_capacity(text_shape, ["one line", "one line"], is_chinese=False)
        assert two_para["used_height_pt"] > one_para["used_height_pt"] * 1.5
