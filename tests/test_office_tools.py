"""Office tools must fail with a clear, actionable message (no raw traceback)
when an optional dependency (openpyxl/python-docx/python-pptx/pdfplumber) is
missing — the fdd-commentary skill relies on this to fall back to the
databook .txt export instead of .xlsx. See AGENT_UPGRADE_NOTES.md section 3.
"""
from __future__ import annotations

import builtins
import json
from pathlib import Path

import pytest

from coding_agent.config import AppConfig, ProviderConfig, RuntimeOptions
from coding_agent.core.runtime import AgentRuntime
from coding_agent.tools.office import _require


def test_require_raises_clear_message_with_pip_extra_and_bare_install() -> None:
    with pytest.raises(RuntimeError) as exc_info:
        _require("no_such_module_xyz", "some-package>=1.0", "excel")
    message = str(exc_info.value)
    assert "some-package>=1.0 is required for this tool" in message
    assert "pip install yucode-agent[excel]" in message
    assert "pip install some-package>=1.0" in message


@pytest.fixture()
def _no_openpyxl(monkeypatch):
    real_import = builtins.__import__

    def fake_import(name, *args, **kwargs):
        if name == "openpyxl" or name.startswith("openpyxl."):
            raise ImportError(f"No module named {name!r}")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", fake_import)


@pytest.fixture()
def runtime(tmp_path: Path) -> AgentRuntime:
    config = AppConfig(
        provider=ProviderConfig(base_url="http://x", api_key="k", model="m"),
        runtime=RuntimeOptions(permission_mode="danger-full-access"),
    )
    return AgentRuntime(tmp_path, config)


def test_read_excel_sheet_fails_cleanly_without_openpyxl(runtime: AgentRuntime, _no_openpyxl) -> None:
    (runtime.workspace_root / "test.xlsx").write_bytes(b"not a real xlsx")
    output = runtime._execute_tool("read_excel_sheet", json.dumps({"path": "test.xlsx"}))

    # Must be the structured tool-error envelope, not a raw Python traceback.
    assert "Traceback" not in output
    assert "File \"" not in output
    payload = json.loads(output)
    assert payload["error_code"] == "tool_error"
    assert "openpyxl" in payload["error"]
    assert "pip install" in payload["error"]


def test_excel_to_json_fails_cleanly_without_openpyxl(runtime: AgentRuntime, _no_openpyxl) -> None:
    (runtime.workspace_root / "test.xlsx").write_bytes(b"not a real xlsx")
    output = runtime._execute_tool("excel_to_json", json.dumps({"path": "test.xlsx"}))
    assert "Traceback" not in output
    payload = json.loads(output)
    assert payload["error_code"] == "tool_error"
    assert "openpyxl" in payload["error"]


# ---------------------------------------------------------------------------
# Shape-level PPTX tools: inspect_pptx_shapes / fill_pptx_shape_text /
# fill_pptx_table (built for the fdd-commentary skill's PPTX export stage).
# Discovery-first design: these tools don't assume any specific template's
# shape-naming convention -- inspect_pptx_shapes reports whatever a given
# template actually contains, and the other two act on shapes found by name.
# ---------------------------------------------------------------------------

@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    """Synthetic template mirroring a house-style deck: a named text
    placeholder and a named table placeholder on one slide."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    textbox.name = "textMainBullets_L"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(4), Inches(1.5))
    table_shape.name = "Table Placeholder"
    path = tmp_path / "template.pptx"
    prs.save(str(path))
    return path


class TestInspectPptxShapes:
    def test_reports_text_and_table_shapes(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        output = runtime._execute_tool("inspect_pptx_shapes", json.dumps({"path": "template.pptx"}))
        payload = json.loads(output)
        shapes = payload["slides"][0]["shapes"]
        names = {s["name"] for s in shapes}
        assert names == {"textMainBullets_L", "Table Placeholder"}

        table_entry = next(s for s in shapes if s["name"] == "Table Placeholder")
        assert table_entry["is_table"] is True
        assert table_entry["rows"] == 2
        assert table_entry["cols"] == 2

        text_entry = next(s for s in shapes if s["name"] == "textMainBullets_L")
        assert text_entry["is_table"] is False
        assert text_entry["width_in"] == 4.0

    def test_missing_dependency_error_is_clean(self, runtime: AgentRuntime, monkeypatch) -> None:
        real_import = builtins.__import__

        def fake_import(name, *args, **kwargs):
            if name == "pptx" or name.startswith("pptx."):
                raise ImportError(f"No module named {name!r}")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        output = runtime._execute_tool("inspect_pptx_shapes", json.dumps({"path": "whatever.pptx"}))
        assert "Traceback" not in output
        payload = json.loads(output)
        assert payload["error_code"] == "tool_error"
        assert "python-pptx" in payload["error"]
        assert "pip install" in payload["error"]


class TestFillPptxShapeText:
    def test_writes_paragraphs_into_named_shape(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        output = runtime._execute_tool("fill_pptx_shape_text", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "textMainBullets_L",
            "paragraphs": ["Cash increased.", "Mainly due to operating inflows."],
        }))
        payload = json.loads(output)
        assert payload["status"] == "ok"
        assert payload["paragraphs_written"] == 2

        from pptx import Presentation
        result = Presentation(str(runtime.workspace_root / "report.pptx"))
        shape = next(s for s in result.slides[0].shapes if s.name == "textMainBullets_L")
        assert "Cash increased." in shape.text_frame.text
        assert "Mainly due to operating inflows." in shape.text_frame.text

    def test_shape_name_lookup_is_case_insensitive(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        output = runtime._execute_tool("fill_pptx_shape_text", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "TEXTMAINBULLETS_L",
            "paragraphs": ["hi"],
        }))
        payload = json.loads(output)
        assert payload["status"] == "ok"

    def test_unknown_shape_name_reports_available_shapes(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        output = runtime._execute_tool("fill_pptx_shape_text", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "doesNotExist",
            "paragraphs": ["hi"],
        }))
        payload = json.loads(output)
        assert payload["error_code"] == "tool_error"
        assert "textMainBullets_L" in payload["error"]

    def test_wrong_shape_type_reports_clear_error(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        output = runtime._execute_tool("fill_pptx_shape_text", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "Table Placeholder",
            "paragraphs": ["hi"],
        }))
        payload = json.loads(output)
        assert payload["error_code"] == "tool_error"
        assert "fill_pptx_table" in payload["error"]

    def test_out_of_range_slide_index_reports_clear_error(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        output = runtime._execute_tool("fill_pptx_shape_text", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 5, "shape_name": "textMainBullets_L",
            "paragraphs": ["hi"],
        }))
        payload = json.loads(output)
        assert payload["error_code"] == "tool_error"
        assert "out of range" in payload["error"]


class TestFillPptxTable:
    def test_rebuilds_table_matching_data_size_exactly(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        """The synthetic template's placeholder table is 2x2 -- 3x3 data must
        NOT get silently cropped to the old size (regression: an earlier
        fill-in-place design did exactly this while claiming success)."""
        output = runtime._execute_tool("fill_pptx_table", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "Table Placeholder",
            "rows": [["Account", "FY24", "FY25"], ["Cash", "10.0", "12.0"], ["AR", "5.0", "6.0"]],
        }))
        payload = json.loads(output)
        assert payload["status"] == "ok"
        assert payload["rows_written"] == 3
        assert payload["cols_written"] == 3

        from pptx import Presentation
        result = Presentation(str(runtime.workspace_root / "report.pptx"))
        table_shape = next(s for s in result.slides[0].shapes if s.has_table)
        table = table_shape.table
        assert len(table.rows) == 3
        assert len(table.columns) == 3
        assert table.cell(0, 0).text == "Account"
        assert table.cell(2, 2).text == "6.0"
        assert table_shape.name == "Table Placeholder"  # name preserved after rebuild

    def test_applies_custom_style_id(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        output = runtime._execute_tool("fill_pptx_table", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "Table Placeholder",
            "rows": [["a", "b"], ["c", "d"]],
            "style_id": "{CUSTOM-TEST-GUID}",
        }))
        payload = json.loads(output)
        assert payload["status"] == "ok"

        from pptx import Presentation
        from pptx.oxml.ns import qn
        result = Presentation(str(runtime.workspace_root / "report.pptx"))
        table_shape = next(s for s in result.slides[0].shapes if s.has_table)
        style_el = table_shape.table._tbl.find(qn("a:tblPr")).find(qn("a:tableStyleId"))
        assert style_el.text == "{CUSTOM-TEST-GUID}"

    def test_progressive_fill_same_output_path(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        """Text and table edits from two separate calls, second reading from
        the first's output, must both survive in the final file."""
        runtime._execute_tool("fill_pptx_shape_text", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "textMainBullets_L",
            "paragraphs": ["Cash commentary."],
        }))
        runtime._execute_tool("fill_pptx_table", json.dumps({
            "path": "report.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "Table Placeholder",
            "rows": [["x", "y"]],
        }))
        from pptx import Presentation
        result = Presentation(str(runtime.workspace_root / "report.pptx"))
        shapes = list(result.slides[0].shapes)
        text_shape = next(s for s in shapes if s.has_text_frame)
        table_shape = next(s for s in shapes if s.has_table)
        assert "Cash commentary." in text_shape.text_frame.text
        assert table_shape.table.cell(0, 0).text == "x"

    def test_rejects_empty_rows(self, runtime: AgentRuntime, sample_pptx: Path) -> None:
        output = runtime._execute_tool("fill_pptx_table", json.dumps({
            "path": "template.pptx", "output_path": "report.pptx",
            "slide_index": 0, "shape_name": "Table Placeholder",
            "rows": [],
        }))
        payload = json.loads(output)
        assert payload["error_code"] == "tool_error"
