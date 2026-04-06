"""MCP server for databook analysis and PPTX generation.

Reads Excel/CSV data, performs basic analysis, generates charts,
and creates PowerPoint presentations from templates or data.

Requires:  openpyxl >= 3.1  (Excel)
Optional:  python-pptx >= 1.0  (PPTX generation)

Run:  python -m coding_agent.plugins.mcp_servers.databook_mcp
"""

from __future__ import annotations

import csv
import io
import json
from pathlib import Path
from typing import Any

from ._protocol import McpStdioServer

_MAX_PREVIEW_ROWS = 50


def _require_openpyxl():
    try:
        import openpyxl
        return openpyxl
    except ImportError:
        raise RuntimeError("openpyxl is required: pip install openpyxl>=3.1")


def _require_pptx():
    try:
        import pptx
        return pptx
    except ImportError:
        raise RuntimeError("python-pptx is required: pip install python-pptx>=1.0")


def _read_csv(path: str, max_rows: int = 500) -> tuple[list[str], list[list[str]]]:
    with open(path, newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        rows = []
        for i, row in enumerate(reader):
            if i > max_rows:
                break
            rows.append(row)
    if not rows:
        return [], []
    return rows[0], rows[1:]


def _read_excel(path: str, sheet: str | None = None, max_rows: int = 500):
    openpyxl = _require_openpyxl()
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    ws = wb[sheet] if sheet else wb.active
    rows = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i > max_rows:
            break
        rows.append([str(c) if c is not None else "" for c in row])
    wb.close()
    if not rows:
        return [], []
    return rows[0], rows[1:]


def handle_read_databook(args: dict[str, Any]) -> Any:
    """Read an Excel or CSV file and return a preview with schema info."""
    path = str(args["path"])
    max_rows = int(args.get("max_rows", _MAX_PREVIEW_ROWS))
    sheet = args.get("sheet")

    ext = Path(path).suffix.lower()
    if ext == ".csv":
        headers, data = _read_csv(path, max_rows)
    elif ext in (".xlsx", ".xls"):
        headers, data = _read_excel(path, sheet, max_rows)
    else:
        raise ValueError(f"Unsupported file format: {ext}. Use .xlsx, .xls, or .csv")

    col_types = []
    for col_idx in range(len(headers)):
        values = [row[col_idx] for row in data[:20] if col_idx < len(row) and row[col_idx]]
        numeric_count = sum(1 for v in values if _is_numeric(v))
        col_type = "numeric" if numeric_count > len(values) * 0.7 else "text"
        col_types.append(col_type)

    return {
        "file": path,
        "format": ext.lstrip("."),
        "headers": headers,
        "column_count": len(headers),
        "row_count": len(data),
        "truncated": len(data) >= max_rows,
        "schema": [
            {"name": h, "inferred_type": col_types[i] if i < len(col_types) else "text"}
            for i, h in enumerate(headers)
        ],
        "preview": data[:10],
    }


def handle_analyze_data(args: dict[str, Any]) -> Any:
    """Basic data analysis: describe, unique counts, or value counts for a column."""
    path = str(args["path"])
    column = args.get("column")
    operation = str(args.get("operation", "describe"))
    sheet = args.get("sheet")
    max_rows = int(args.get("max_rows", 5000))

    ext = Path(path).suffix.lower()
    if ext == ".csv":
        headers, data = _read_csv(path, max_rows)
    else:
        headers, data = _read_excel(path, sheet, max_rows)

    if not headers:
        return {"error": "Empty dataset"}

    if column and column not in headers:
        return {"error": f"Column '{column}' not found. Available: {headers}"}

    if operation == "describe":
        result: dict[str, Any] = {
            "row_count": len(data),
            "column_count": len(headers),
            "columns": {},
        }
        for i, h in enumerate(headers):
            values = [row[i] for row in data if i < len(row) and row[i]]
            nums = [float(v) for v in values if _is_numeric(v)]
            col_info: dict[str, Any] = {
                "non_empty": len(values),
                "unique": len(set(values)),
            }
            if nums:
                col_info["min"] = min(nums)
                col_info["max"] = max(nums)
                col_info["mean"] = round(sum(nums) / len(nums), 4)
            result["columns"][h] = col_info
        return result

    if operation == "value_counts" and column:
        col_idx = headers.index(column)
        counts: dict[str, int] = {}
        for row in data:
            if col_idx < len(row):
                val = row[col_idx]
                counts[val] = counts.get(val, 0) + 1
        sorted_counts = sorted(counts.items(), key=lambda x: x[1], reverse=True)[:50]
        return {"column": column, "value_counts": dict(sorted_counts)}

    if operation == "unique" and column:
        col_idx = headers.index(column)
        unique_vals = sorted(set(row[col_idx] for row in data if col_idx < len(row)))[:100]
        return {"column": column, "unique_count": len(unique_vals), "values": unique_vals}

    return {"error": f"Unknown operation: {operation}"}


def handle_create_pptx(args: dict[str, Any]) -> Any:
    """Create a PPTX from structured slide data, optionally using a template."""
    pptx = _require_pptx()
    output_path = str(args["output_path"])
    slides_data = args.get("slides", [])
    template_path = args.get("template_path")

    if template_path:
        prs = pptx.Presentation(template_path)
    else:
        prs = pptx.Presentation()

    for slide_info in slides_data:
        layout_idx = int(slide_info.get("layout", 1))
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 1
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        title = slide_info.get("title", "")
        if title and slide.shapes.title:
            slide.shapes.title.text = title

        body = slide_info.get("body", "")
        bullets = slide_info.get("bullets", [])
        if slide.placeholders and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            if body:
                tf.text = body
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = str(bullet)

        notes = slide_info.get("notes", "")
        if notes and hasattr(slide, "notes_slide"):
            slide.notes_slide.notes_text_frame.text = notes

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return {"created": output_path, "slide_count": len(slides_data)}


def handle_data_to_pptx(args: dict[str, Any]) -> Any:
    """Read a databook and create a summary PPTX with the data."""
    pptx = _require_pptx()
    data_path = str(args["data_path"])
    output_path = str(args["output_path"])
    title = str(args.get("title", Path(data_path).stem))
    sheet = args.get("sheet")
    max_rows = int(args.get("max_rows", 100))

    ext = Path(data_path).suffix.lower()
    if ext == ".csv":
        headers, data = _read_csv(data_path, max_rows)
    else:
        headers, data = _read_excel(data_path, sheet, max_rows)

    prs = pptx.Presentation()

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = title
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = f"{len(data)} rows x {len(headers)} columns"

    schema_layout = prs.slide_layouts[1]
    schema_slide = prs.slides.add_slide(schema_layout)
    if schema_slide.shapes.title:
        schema_slide.shapes.title.text = "Data Schema"
    if schema_slide.placeholders and len(schema_slide.placeholders) > 1:
        tf = schema_slide.placeholders[1].text_frame
        tf.text = f"Columns ({len(headers)}):"
        for h in headers[:20]:
            p = tf.add_paragraph()
            p.text = f"  {h}"

    chunk_size = 15
    for chunk_start in range(0, min(len(data), 60), chunk_size):
        chunk_end = min(chunk_start + chunk_size, len(data))
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = f"Data (rows {chunk_start + 1}-{chunk_end})"
        if slide.placeholders and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            for row in data[chunk_start:chunk_end]:
                p = tf.add_paragraph()
                p.text = " | ".join(row[:8])

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return {"created": output_path, "slide_count": len(prs.slides)}


def _is_numeric(value: str) -> bool:
    try:
        float(value.replace(",", ""))
        return True
    except (ValueError, AttributeError):
        return False


def main():
    server = McpStdioServer("yucode-databook")

    server.register_tool("read_databook", "Read an Excel or CSV file and return a preview with schema inference.", {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Path to .xlsx, .xls, or .csv file"},
            "sheet": {"type": "string", "description": "Sheet name for Excel (default: active sheet)"},
            "max_rows": {"type": "integer", "description": "Max rows to read (default 50)"},
        },
        "required": ["path"],
    }, handle_read_databook)

    server.register_tool("analyze_data", "Analyze data from Excel/CSV: describe, value_counts, or unique.", {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Path to data file"},
            "sheet": {"type": "string", "description": "Sheet name (Excel only)"},
            "column": {"type": "string", "description": "Column name to analyze"},
            "operation": {"type": "string", "description": "Operation: describe (default), value_counts, unique"},
            "max_rows": {"type": "integer", "description": "Max rows to load (default 5000)"},
        },
        "required": ["path"],
    }, handle_analyze_data)

    server.register_tool("create_pptx", "Create a PowerPoint file from structured slide data.", {
        "type": "object",
        "properties": {
            "output_path": {"type": "string", "description": "Output .pptx file path"},
            "template_path": {"type": "string", "description": "Optional template .pptx to use as base"},
            "slides": {
                "type": "array",
                "description": "List of slides. Each: {title, body, bullets: [...], layout: int, notes: str}",
                "items": {"type": "object"},
            },
        },
        "required": ["output_path", "slides"],
    }, handle_create_pptx)

    server.register_tool("data_to_pptx", "Read a databook and create a summary PPTX presentation.", {
        "type": "object",
        "properties": {
            "data_path": {"type": "string", "description": "Path to .xlsx or .csv data file"},
            "output_path": {"type": "string", "description": "Output .pptx file path"},
            "title": {"type": "string", "description": "Presentation title (default: filename)"},
            "sheet": {"type": "string", "description": "Sheet name (Excel only)"},
            "max_rows": {"type": "integer", "description": "Max data rows (default 100)"},
        },
        "required": ["data_path", "output_path"],
    }, handle_data_to_pptx)

    server.serve()


if __name__ == "__main__":
    main()
