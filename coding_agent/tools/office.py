"""Built-in Office file tools -- Word (.docx), Excel (.xlsx), PowerPoint (.pptx).

These are direct tools (no MCP server needed). Each gracefully degrades
if the optional dependency is missing, telling the model how to install it.
"""

from __future__ import annotations

import base64
import contextlib
import json
import mimetypes
from typing import TYPE_CHECKING, Any

from . import RiskLevel, ToolDefinition, ToolSpec

if TYPE_CHECKING:
    from . import ToolRegistry


_IMAGE_MAX_BASE64_BYTES = 1_500_000  # ~1 MB cap on raw bytes before base64-encoding
_IMAGE_SUFFIXES = (".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff")


def office_tools(registry: ToolRegistry) -> list[ToolDefinition]:
    return [
        # ---- Excel ----
        ToolDefinition(
            ToolSpec(
                "read_excel_sheet",
                "Read an Excel (.xlsx) sheet as headers + rows. Returns JSON with 'headers' and 'rows'.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .xlsx file."},
                        "sheet": {"type": "string", "description": "Sheet name (default: active sheet)."},
                        "max_rows": {"type": "integer", "description": "Max data rows to return (default 500)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _read_excel_sheet(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "list_excel_sheets",
                "List visible sheet names in an Excel (.xlsx) file. Hidden sheets are excluded.",
                {
                    "type": "object",
                    "properties": {"path": {"type": "string", "description": "Path to .xlsx file."}},
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _list_excel_sheets(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "inspect_excel_sheets",
                "Scan all visible sheets in an Excel file. Returns each sheet's name, "
                "estimated row count, column headers, and sample rows. "
                "Call this before read_excel_sheet when you need to identify which sheet "
                "contains the data the user is looking for.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .xlsx file."},
                        "sample_rows": {"type": "integer", "description": "Data rows to sample per sheet (default 3)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _inspect_excel_sheets(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "write_excel_cell",
                "Write a value to a cell in an Excel (.xlsx) file.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "sheet": {"type": "string", "description": "Sheet name (default: active)."},
                        "cell": {"type": "string", "description": "Cell reference like A1, B2."},
                        "value": {"type": "string", "description": "Value to write."},
                    },
                    "required": ["path", "cell", "value"],
                },
                "workspace-write",
                RiskLevel.MEDIUM,
            ),
            lambda args: _write_excel_cell(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "excel_to_json",
                "Convert an Excel sheet to a list of JSON records (using first row as headers).",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "sheet": {"type": "string"},
                        "max_rows": {"type": "integer", "description": "Max rows (default 500)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _excel_to_json(registry, args),
        ),
        # ---- Word ----
        ToolDefinition(
            ToolSpec(
                "read_word_text",
                "Read the full text of a Word (.docx) file.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .docx file."},
                        "max_chars": {"type": "integer", "description": "Max characters (default 50000)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _read_word_text(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "read_word_paragraphs",
                "Read paragraphs from a Word (.docx) file with style info.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "max_paragraphs": {"type": "integer", "description": "Max paragraphs (default 1000)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _read_word_paragraphs(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "write_word",
                "Create a new Word (.docx) file with paragraphs.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "paragraphs": {
                            "type": "array",
                            "description": "List of paragraphs. Each can be a string or {text, style}.",
                            "items": {},
                        },
                    },
                    "required": ["path", "paragraphs"],
                },
                "workspace-write",
                RiskLevel.MEDIUM,
            ),
            lambda args: _write_word(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "append_word",
                "Append a paragraph to an existing Word (.docx) file.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "text": {"type": "string"},
                        "style": {"type": "string", "description": "Paragraph style (default Normal)."},
                    },
                    "required": ["path", "text"],
                },
                "workspace-write",
                RiskLevel.MEDIUM,
            ),
            lambda args: _append_word(registry, args),
        ),
        # ---- PowerPoint ----
        ToolDefinition(
            ToolSpec(
                "read_pptx",
                "Read text content from a PowerPoint (.pptx) file. Returns slide-by-slide text.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pptx file."},
                        "max_slides": {"type": "integer", "description": "Max slides to read (default 100)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _read_pptx(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "write_pptx",
                "Create a new PowerPoint (.pptx) file with slides. Each slide has a title and body.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "slides": {
                            "type": "array",
                            "description": "List of slides, each {title, body} or {title, bullets: [...]}.",
                            "items": {},
                        },
                    },
                    "required": ["path", "slides"],
                },
                "workspace-write",
                RiskLevel.MEDIUM,
            ),
            lambda args: _write_pptx(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "write_pptx_from_template",
                "Create a PowerPoint (.pptx) from an existing template. "
                "Replaces placeholders in the template and adds slides with data.",
                {
                    "type": "object",
                    "properties": {
                        "template_path": {"type": "string", "description": "Path to template .pptx file."},
                        "output_path": {"type": "string", "description": "Output .pptx file path."},
                        "slides": {
                            "type": "array",
                            "description": "List of slides: each {title, body, bullets: [...], notes: str}.",
                            "items": {},
                        },
                    },
                    "required": ["template_path", "output_path", "slides"],
                },
                "workspace-write",
                RiskLevel.MEDIUM,
            ),
            lambda args: _write_pptx_from_template(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "inspect_pptx_shapes",
                "Inspect a PowerPoint (.pptx) file's actual shape structure, per slide: each "
                "shape's name, type, position, and (for tables) row/column count with a cell "
                "preview, or (for text frames) a text preview. Call this before "
                "fill_pptx_shape_text/fill_pptx_table on a template whose shape names you don't "
                "already know — house-style templates typically name their content "
                "placeholders (e.g. a shape called 'textMainBullets' or 'Table Placeholder').",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pptx file."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _inspect_pptx_shapes(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "fill_pptx_shape_text",
                "Write text into a specific named shape (text box or placeholder) on a specific "
                "slide of a .pptx file. Matches shape name case-insensitively. Use "
                "inspect_pptx_shapes first to find the right slide_index/shape_name. Pass the "
                "same path as output_path to keep progressively filling the same in-progress file. "
                "Applies the same fixed formatting estimate_pptx_text_capacity assumes (9pt "
                "Arial/Microsoft YaHei, single spacing, 3pt paragraph gap) so the two tools never "
                "disagree, and sets the shape's autofit accordingly: noAutofit when the text fits "
                "(exact size, no surprise PowerPoint shrink), or a bounded shrink-to-fit when it "
                "doesn't (visible, not silently clipped).",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pptx file to read."},
                        "output_path": {"type": "string", "description": "Path to save the result to."},
                        "slide_index": {"type": "integer", "description": "0-based slide index."},
                        "shape_name": {"type": "string", "description": "Target shape's name (case-insensitive)."},
                        "paragraphs": {
                            "type": "array",
                            "description": "List of paragraph strings; each becomes one paragraph/bullet.",
                            "items": {"type": "string"},
                        },
                        "is_chinese": {
                            "type": "boolean",
                            "description": "Optional override for font selection; auto-detected (>30% CJK characters) if omitted.",
                        },
                    },
                    "required": ["path", "output_path", "slide_index", "shape_name", "paragraphs"],
                },
                "workspace-write",
                RiskLevel.MEDIUM,
            ),
            lambda args: _fill_pptx_shape_text(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "fill_pptx_table",
                "Write row data into a named table (or table placeholder) shape on a specific "
                "slide of a .pptx file. The target shape is always replaced with a freshly built "
                "table sized exactly to the given data, at the same position — this never "
                "silently crops data to an existing table's size. Optionally pass style_id to set "
                "a table style GUID. Matches shape name case-insensitively. Use "
                "inspect_pptx_shapes first to find slide_index/shape_name.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pptx file to read."},
                        "output_path": {"type": "string", "description": "Path to save the result to."},
                        "slide_index": {"type": "integer", "description": "0-based slide index."},
                        "shape_name": {"type": "string", "description": "Target shape's name (case-insensitive)."},
                        "rows": {
                            "type": "array",
                            "description": "Row data: list of rows, each a list of cell strings.",
                            "items": {"type": "array", "items": {"type": "string"}},
                        },
                        "style_id": {
                            "type": "string",
                            "description": "Optional table style GUID (e.g. '{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}') to apply.",
                        },
                    },
                    "required": ["path", "output_path", "slide_index", "shape_name", "rows"],
                },
                "workspace-write",
                RiskLevel.MEDIUM,
            ),
            lambda args: _fill_pptx_table(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "estimate_pptx_text_capacity",
                "Estimate whether commentary text fits a named shape's slot on a .pptx slide, "
                "using real font-glyph measurement (not a rough chars-per-inch guess). Pass the "
                "SAME paragraphs list you're about to write with fill_pptx_shape_text. Returns "
                "fill_ratio and remaining_lines_estimate -- use these to decide whether to add "
                "more genuinely-supported detail (a slot at fill_ratio 0.2-0.3 has room to spare) "
                "or trim (fill_ratio > 1.0 means it overflows). Call inspect_pptx_shapes first to "
                "find slide_index/shape_name.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pptx file."},
                        "slide_index": {"type": "integer", "description": "0-based slide index."},
                        "shape_name": {"type": "string", "description": "Target shape's name (case-insensitive)."},
                        "paragraphs": {
                            "type": "array",
                            "description": "Candidate paragraph strings, same shape as fill_pptx_shape_text's `paragraphs`.",
                            "items": {"type": "string"},
                        },
                        "is_chinese": {
                            "type": "boolean",
                            "description": "Optional override; auto-detected (>30% CJK characters) if omitted.",
                        },
                    },
                    "required": ["path", "slide_index", "shape_name", "paragraphs"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _estimate_pptx_text_capacity(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "locate_databook_stage_columns",
                "Scan a financial databook sheet's header rows for reporting 'basis/stage' "
                "labels (Indicative adjusted, Mgt acc, Audited, Audit adjustment, Indicative "
                "adjustment -- including common Traditional/Simplified Chinese and typo "
                "variants) and return exactly which column each canonical stage occupies, "
                "with its associated period date. Call this BEFORE reading account figures "
                "from a databook, instead of visually picking which of many similarly-dated "
                "column groups is the required one (e.g. 'use only the Indicative adjusted "
                "columns') -- real databooks routinely place 10-30+ date columns side by "
                "side across 2-5 basis groups with only a text label distinguishing them, no "
                "cell merging or other structural hint. Returns found=false if the sheet "
                "doesn't have a recognizable multi-basis header.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .xlsx file."},
                        "sheet": {"type": "string", "description": "Sheet name (default: active sheet)."},
                        "max_scan_rows": {"type": "integer", "description": "Max rows to scan for the basis/header row (default 20)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _locate_databook_stage_columns(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "read_excel_preview",
                "Read an Excel file with schema inference. Returns headers, types, and a preview of the data.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .xlsx file."},
                        "sheet": {"type": "string", "description": "Sheet name (default: active sheet)."},
                        "max_rows": {"type": "integer", "description": "Max preview rows (default 20)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _read_excel_preview(registry, args),
        ),
        # ---- Image ----
        ToolDefinition(
            ToolSpec(
                "image_read",
                "Inspect an image (.png/.jpg/.jpeg/.webp/.gif/.bmp/.tiff) and return "
                "metadata (path, mime_type, width, height, mode, size_bytes). "
                "Pass include_base64=true to also get a base64 data URL — useful when "
                "the provider supports multimodal/vision input. Default omits base64 "
                "to keep the response cheap.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to image file."},
                        "include_base64": {"type": "boolean", "description": "If true, include a base64 data URL in the response."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _image_read(registry, args),
        ),
        # ---- PDF ----
        ToolDefinition(
            ToolSpec(
                "read_pdf_text",
                "Extract text from a PDF file. Returns page-by-page text. Pass layout=true to "
                "preserve the printed column positions with whitespace (pdftotext -layout "
                "equivalent) — use that mode when reading financial statements or any page "
                "where column alignment matters.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pdf file."},
                        "max_pages": {"type": "integer", "description": "Max pages to read (default 50)."},
                        "layout": {"type": "boolean", "description": "Preserve printed column layout with whitespace (default false)."},
                    },
                    "required": ["path"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _read_pdf_text(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "extract_pdf_words",
                "All words on one PDF page with their x/y coordinates and x-centres. Low-level "
                "building block for coordinate-based table extraction when "
                "extract_pdf_period_table's automatic banding isn't enough (unusual layouts, "
                "matrix tables with text column headers).",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pdf file."},
                        "page": {"type": "integer", "description": "1-based page number."},
                    },
                    "required": ["path", "page"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _extract_pdf_words(registry, args),
        ),
        ToolDefinition(
            ToolSpec(
                "extract_pdf_period_table",
                "Column-safe extraction of a period-columned financial table from one PDF page. "
                "Detects the period headers (e.g. 2025 / 2024) by x-centre, builds a horizontal "
                "band per column, groups words into visual rows by y, and assigns each value to "
                "the period whose band its x-centre falls in — NEVER by token/reading order. "
                "Printed dashes count as 0 only inside a band; bracketed figures come back as "
                "real negatives; wrapped label lines are stitched onto their valued row; "
                "right-hand translated labels (bilingual statements) are kept separately and "
                "never counted as values. Use this for EVERY bilingual side-by-side note table. "
                "Returns found=false (with a reason) instead of guessing when the page has no "
                "text layer or no period header — flag those pages, do not invent columns.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pdf file."},
                        "page": {"type": "integer", "description": "1-based page number."},
                        "periods": {
                            "type": "array", "items": {"type": "string"},
                            "description": "Expected period headers as printed (e.g. [\"2025\", \"2024\"]). Omit to auto-detect year-like headers.",
                        },
                    },
                    "required": ["path", "page"],
                },
                "read-only",
                RiskLevel.LOW,
            ),
            lambda args: _extract_pdf_period_table(registry, args),
        ),
    ]


# ---- Excel handlers ---------------------------------------------------------

def _require(module: str, package: str, pip_extra: str) -> Any:
    try:
        return __import__(module)
    except ImportError as exc:
        raise RuntimeError(
            f"{package} is required for this tool. "
            f"Install with: pip install yucode-agent[{pip_extra}]  or  pip install {package}"
        ) from exc


def _get_visibility_state(ws: Any) -> tuple[set[int], set[int]]:
    """Return (hidden_row_numbers_1based, hidden_col_indices_1based).

    Requires a fully-loaded Worksheet (not ReadOnlyWorksheet) so that
    row_dimensions and column_dimensions are available.
    """
    from openpyxl.utils import column_index_from_string
    hidden_rows: set[int] = set()
    hidden_cols: set[int] = set()
    for row_num, rd in ws.row_dimensions.items():
        if rd.hidden:
            hidden_rows.add(int(row_num))
    for col_letter, cd in ws.column_dimensions.items():
        if cd.hidden:
            with contextlib.suppress(Exception):
                hidden_cols.add(column_index_from_string(col_letter))
    return hidden_rows, hidden_cols


def _read_excel_sheet(registry: ToolRegistry, args: dict[str, Any]) -> str:
    openpyxl = _require("openpyxl", "openpyxl>=3.1", "excel")
    path = registry._resolve_path(str(args["path"]))
    wb = openpyxl.load_workbook(str(path), data_only=True)
    sheet_name = args.get("sheet")
    ws = wb[sheet_name] if sheet_name else wb.active
    max_rows = int(args.get("max_rows", 500))
    hidden_rows, hidden_cols = _get_visibility_state(ws)
    rows: list[list[str]] = []
    skipped = 0
    for row_cells in ws.iter_rows():
        if not row_cells:
            continue
        row_num = row_cells[0].row
        if row_num in hidden_rows:
            skipped += 1
            continue
        if len(rows) >= max_rows + 1:
            break
        rows.append([
            str(c.value) if c.value is not None else ""
            for c in row_cells if c.column not in hidden_cols
        ])
    wb.close()
    if not rows:
        return json.dumps({"headers": [], "rows": []})
    result: dict[str, Any] = {"headers": rows[0], "rows": rows[1:]}
    if skipped or hidden_cols:
        result["note"] = f"Skipped {skipped} hidden row(s) and {len(hidden_cols)} hidden column(s)."
    return json.dumps(result)


def _list_excel_sheets(registry: ToolRegistry, args: dict[str, Any]) -> str:
    openpyxl = _require("openpyxl", "openpyxl>=3.1", "excel")
    path = registry._resolve_path(str(args["path"]))
    wb = openpyxl.load_workbook(str(path), read_only=True)
    visible = [ws.title for ws in wb.worksheets if ws.sheet_state == "visible"]
    hidden_count = sum(1 for ws in wb.worksheets if ws.sheet_state != "visible")
    wb.close()
    result: dict[str, Any] = {"sheets": visible}
    if hidden_count:
        result["hidden_sheets_count"] = hidden_count
    return json.dumps(result)


def _write_excel_cell(registry: ToolRegistry, args: dict[str, Any]) -> str:
    openpyxl = _require("openpyxl", "openpyxl>=3.1", "excel")
    path = registry._resolve_path(str(args["path"]))
    wb = openpyxl.load_workbook(str(path))
    sheet_name = args.get("sheet")
    ws = wb[sheet_name] if sheet_name else wb.active
    cell = args["cell"]
    ws[cell] = args["value"]
    wb.save(str(path))
    wb.close()
    return f"Wrote '{args['value']}' to {cell} in {path.name}"


def _excel_to_json(registry: ToolRegistry, args: dict[str, Any]) -> str:
    openpyxl = _require("openpyxl", "openpyxl>=3.1", "excel")
    path = registry._resolve_path(str(args["path"]))
    wb = openpyxl.load_workbook(str(path), data_only=True)
    sheet_name = args.get("sheet")
    ws = wb[sheet_name] if sheet_name else wb.active
    max_rows = int(args.get("max_rows", 500))
    hidden_rows, hidden_cols = _get_visibility_state(ws)
    rows: list[list[str]] = []
    for row_cells in ws.iter_rows():
        if not row_cells:
            continue
        if row_cells[0].row in hidden_rows:
            continue
        if len(rows) >= max_rows + 1:
            break
        rows.append([
            str(c.value) if c.value is not None else ""
            for c in row_cells if c.column not in hidden_cols
        ])
    wb.close()
    if not rows:
        return "[]"
    headers = [str(h) if h else f"col_{i}" for i, h in enumerate(rows[0])]
    result = []
    for row in rows[1:]:
        result.append({
            headers[j] if j < len(headers) else f"col_{j}": v
            for j, v in enumerate(row)
        })
    return json.dumps(result, indent=2)


# ---- Word handlers ----------------------------------------------------------

def _read_word_text(registry: ToolRegistry, args: dict[str, Any]) -> str:
    docx = _require("docx", "python-docx>=1.1", "word")
    path = registry._resolve_path(str(args["path"]))
    doc = docx.Document(str(path))
    text = "\n".join(p.text for p in doc.paragraphs)
    max_chars = int(args.get("max_chars", 50000))
    if len(text) > max_chars:
        text = text[:max_chars] + f"\n... truncated at {max_chars} chars"
    return text


def _read_word_paragraphs(registry: ToolRegistry, args: dict[str, Any]) -> str:
    docx = _require("docx", "python-docx>=1.1", "word")
    path = registry._resolve_path(str(args["path"]))
    doc = docx.Document(str(path))
    max_paras = int(args.get("max_paragraphs", 1000))
    paragraphs = []
    for i, para in enumerate(doc.paragraphs):
        if i >= max_paras:
            break
        paragraphs.append({"index": i, "style": para.style.name, "text": para.text})
    return json.dumps({"paragraphs": paragraphs, "total": len(doc.paragraphs)})


def _write_word(registry: ToolRegistry, args: dict[str, Any]) -> str:
    docx = _require("docx", "python-docx>=1.1", "word")
    path = registry._resolve_path(str(args["path"]))
    path.parent.mkdir(parents=True, exist_ok=True)
    doc = docx.Document()
    for para in args.get("paragraphs", []):
        text = para if isinstance(para, str) else str(para.get("text", ""))
        style = "Normal" if isinstance(para, str) else str(para.get("style", "Normal"))
        doc.add_paragraph(text, style=style)
    doc.save(str(path))
    return f"Created {path}"


def _append_word(registry: ToolRegistry, args: dict[str, Any]) -> str:
    docx = _require("docx", "python-docx>=1.1", "word")
    path = registry._resolve_path(str(args["path"]))
    doc = docx.Document(str(path))
    style = args.get("style", "Normal")
    doc.add_paragraph(args["text"], style=style)
    doc.save(str(path))
    return f"Appended paragraph to {path}"


# ---- PowerPoint handlers ----------------------------------------------------

def _read_pptx(registry: ToolRegistry, args: dict[str, Any]) -> str:
    pptx = _require("pptx", "python-pptx>=1.0", "pptx")
    path = registry._resolve_path(str(args["path"]))
    prs = pptx.Presentation(str(path))
    max_slides = int(args.get("max_slides", 100))
    slides = []
    for i, slide in enumerate(prs.slides):
        if i >= max_slides:
            break
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides.append({"slide": i + 1, "text": "\n".join(texts)})
    return json.dumps({"slides": slides, "total": len(prs.slides)})


def _write_pptx(registry: ToolRegistry, args: dict[str, Any]) -> str:
    pptx = _require("pptx", "python-pptx>=1.0", "pptx")
    path = registry._resolve_path(str(args["path"]))
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = pptx.Presentation()
    for slide_data in args.get("slides", []):
        if isinstance(slide_data, str):
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_data
            continue
        title = slide_data.get("title", "")
        body = slide_data.get("body", "")
        bullets = slide_data.get("bullets", [])
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if title and slide.shapes.title:
            slide.shapes.title.text = title
        if slide.placeholders and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            if body:
                tf.text = body
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = str(bullet)
    prs.save(str(path))
    return f"Created {path} with {len(args.get('slides', []))} slides"


# ---- Template PPTX handler --------------------------------------------------

def _write_pptx_from_template(registry: ToolRegistry, args: dict[str, Any]) -> str:
    pptx = _require("pptx", "python-pptx>=1.0", "pptx")
    template_path = registry._resolve_path(str(args["template_path"]))
    output_path = registry._resolve_path(str(args["output_path"]))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = pptx.Presentation(str(template_path))
    for slide_data in args.get("slides", []):
        layout_idx = int(slide_data.get("layout", 1))
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 1
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        title = slide_data.get("title", "")
        if title and slide.shapes.title:
            slide.shapes.title.text = title
        body = slide_data.get("body", "")
        bullets = slide_data.get("bullets", [])
        if slide.placeholders and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            if body:
                tf.text = body
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = str(bullet)
        notes = slide_data.get("notes", "")
        if notes and hasattr(slide, "notes_slide"):
            slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(output_path))
    return f"Created {output_path} from template with {len(args.get('slides', []))} slides"


# ---- PPTX shape-level inspection / filling ----------------------------------
# Discovery-first design: rather than assuming a fixed template's shape names
# (e.g. a specific firm's house style), inspect_pptx_shapes lets the agent
# learn ANY given template's actual structure at runtime, then fill_pptx_shape_text
# / fill_pptx_table act on whatever names that discovery turns up.

_EMU_PER_INCH = 914400


def _shape_dims_in(shape: Any) -> dict[str, float | None]:
    def _to_inches(value: Any) -> float | None:
        return round(value / _EMU_PER_INCH, 2) if value is not None else None
    return {
        "left_in": _to_inches(shape.left),
        "top_in": _to_inches(shape.top),
        "width_in": _to_inches(shape.width),
        "height_in": _to_inches(shape.height),
    }


def _find_shape_by_name(slide: Any, shape_name: str) -> Any:
    target_name = shape_name.strip().lower()
    for shape in slide.shapes:
        if shape.name.strip().lower() == target_name:
            return shape
    available = [shape.name for shape in slide.shapes]
    raise ValueError(f"No shape named '{shape_name}' on this slide. Available shapes: {available}")


def _resolve_slide(prs: Any, slide_index: int) -> Any:
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"slide_index {slide_index} out of range (presentation has {len(prs.slides)} slides).")
    return prs.slides[slide_index]


def _inspect_pptx_shapes(registry: ToolRegistry, args: dict[str, Any]) -> str:
    pptx = _require("pptx", "python-pptx>=1.0", "pptx")
    path = registry._resolve_path(str(args["path"]))
    prs = pptx.Presentation(str(path))
    slides_info = []
    for slide_index, slide in enumerate(prs.slides):
        shapes_info = []
        for shape in slide.shapes:
            entry: dict[str, Any] = {"name": shape.name, "shape_type": str(shape.shape_type)}
            entry.update(_shape_dims_in(shape))
            if shape.has_table:
                table = shape.table
                entry["is_table"] = True
                entry["rows"] = len(table.rows)
                entry["cols"] = len(table.columns)
                entry["preview"] = [
                    [cell.text for cell in row.cells] for row in list(table.rows)[:3]
                ]
            else:
                entry["is_table"] = False
                if shape.has_text_frame:
                    entry["text_preview"] = shape.text_frame.text[:150]
            shapes_info.append(entry)
        slides_info.append({"slide_index": slide_index, "shapes": shapes_info})
    return json.dumps({"slides": slides_info}, indent=2)


_AUTOFIT_TAGS = ("a:spAutoFit", "a:normAutofit", "a:noAutofit")
_BOUNDED_AUTOFIT_MIN_SCALE = 0.70


def _clear_autofit(body_pr: Any) -> None:
    from pptx.oxml.ns import qn
    for tag in _AUTOFIT_TAGS:
        for child in body_pr.findall(qn(tag)):
            body_pr.remove(child)


def _force_no_autofit(text_frame: Any) -> None:
    """Set bodyPr autofit to <a:noAutofit/> so text renders at exactly the
    size just applied. A freshly-added or template placeholder shape often
    ships with <a:spAutoFit/> (resize shape to text) or a stale
    <a:normAutofit fontScale=".."/> computed by PowerPoint for whatever text
    was there before -- either would silently override the size for newly
    written content that estimate_pptx_text_capacity never accounted for."""
    from pptx.oxml import parse_xml
    body_pr = text_frame._txBody.bodyPr
    _clear_autofit(body_pr)
    body_pr.append(parse_xml(
        '<a:noAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    ))


def _apply_bounded_autofit(text_frame: Any, font_scale: float) -> None:
    """Set a real <a:normAutofit> so PowerPoint visibly shrinks text that
    doesn't fit at the fixed size, instead of silently clipping it at the
    shape edge. Bounded at _BOUNDED_AUTOFIT_MIN_SCALE so extreme overflow
    still clips rather than rendering illegibly small text."""
    from pptx.oxml import parse_xml
    scale = max(_BOUNDED_AUTOFIT_MIN_SCALE, min(1.0, float(font_scale)))
    font_pct = int(round(scale * 100000))
    line_pct = int(round(max(0.0, (1.0 - scale) * 0.5) * 100000))
    body_pr = text_frame._txBody.bodyPr
    _clear_autofit(body_pr)
    body_pr.append(parse_xml(
        f'<a:normAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'fontScale="{font_pct}" lnSpcReduction="{line_pct}"/>'
    ))


def _fill_pptx_shape_text(registry: ToolRegistry, args: dict[str, Any]) -> str:
    pptx = _require("pptx", "python-pptx>=1.0", "pptx")
    _require("PIL", "Pillow>=9.0", "pptx")
    from pptx.util import Pt

    from . import text_metrics

    path = registry._resolve_path(str(args["path"]))
    output_path = registry._resolve_path(str(args["output_path"]))
    slide_index = int(args["slide_index"])
    shape_name = str(args["shape_name"])
    paragraphs = [str(p) for p in args.get("paragraphs", [])]
    is_chinese = args.get("is_chinese")
    if is_chinese is not None:
        is_chinese = bool(is_chinese)

    prs = pptx.Presentation(str(path))
    slide = _resolve_slide(prs, slide_index)
    shape = _find_shape_by_name(slide, shape_name)
    if not shape.has_text_frame:
        raise ValueError(
            f"Shape '{shape.name}' on slide {slide_index} has no text frame "
            f"(is_table={shape.has_table}) — use fill_pptx_table instead."
        )

    chinese = (
        text_metrics.is_predominantly_chinese("\n".join(paragraphs))
        if is_chinese is None else is_chinese
    )
    font_name = "Microsoft YaHei" if chinese else "Arial"

    # Same fixed 9pt / single-spacing / 3pt-gap formatting that
    # estimate_pptx_text_capacity assumes when predicting fit -- without
    # this, the two tools silently disagree: capacity says "fits" for
    # formatting that fill_pptx_shape_text never actually applies.
    tf = shape.text_frame
    tf.clear()
    for i, para_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para_text
        p.line_spacing = text_metrics.ESTIMATE_LINE_SPACING
        p.space_before = Pt(0)
        p.space_after = Pt(text_metrics.ESTIMATE_PARA_GAP_PT)
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(text_metrics.ESTIMATE_FONT_SIZE_PT)

    try:
        capacity = text_metrics.estimate_capacity(shape, paragraphs, is_chinese=chinese)
    except FileNotFoundError:
        capacity = None
    if capacity is None or capacity["fits"]:
        _force_no_autofit(tf)
        autofit = "none"
    else:
        used = capacity["used_height_pt"]
        font_scale = capacity["box_height_pt"] / used if used > 0 else 1.0
        _apply_bounded_autofit(tf, font_scale)
        autofit = "bounded_shrink"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return json.dumps({
        "status": "ok", "slide_index": slide_index, "shape_name": shape.name,
        "paragraphs_written": len(paragraphs), "output_path": str(output_path),
        "font_name": font_name, "font_size_pt": text_metrics.ESTIMATE_FONT_SIZE_PT,
        "autofit": autofit,
    }, indent=2)


def _set_table_style_id(table: Any, style_id: str) -> None:
    from pptx.oxml.ns import qn
    tbl_pr = table._tbl.find(qn("a:tblPr"))
    if tbl_pr is None:
        return
    style_el = tbl_pr.find(qn("a:tableStyleId"))
    if style_el is not None:
        style_el.text = style_id


def _fill_pptx_table(registry: ToolRegistry, args: dict[str, Any]) -> str:
    pptx = _require("pptx", "python-pptx>=1.0", "pptx")
    path = registry._resolve_path(str(args["path"]))
    output_path = registry._resolve_path(str(args["output_path"]))
    slide_index = int(args["slide_index"])
    shape_name = str(args["shape_name"])
    rows_data = [[str(cell) for cell in row] for row in args["rows"]]
    style_id = str(args.get("style_id", ""))

    n_rows = len(rows_data)
    n_cols = max((len(row) for row in rows_data), default=0)
    if n_rows == 0 or n_cols == 0:
        raise ValueError("`rows` must be a non-empty list of non-empty lists.")

    prs = pptx.Presentation(str(path))
    slide = _resolve_slide(prs, slide_index)
    shape = _find_shape_by_name(slide, shape_name)

    # Always delete-and-rebuild rather than filling an existing table in place:
    # an existing table's row/col count rarely matches fresh data exactly, and
    # silently cropping to whatever size was already there is a worse surprise
    # than a clean rebuild — this also matches the reference project's own
    # approach (it never reuses a placeholder's table, always rebuilds it).
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    original_name = shape.name
    element = shape._element
    element.getparent().remove(element)
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    graphic_frame.name = original_name
    new_table = graphic_frame.table
    for r_idx, row_values in enumerate(rows_data):
        for c_idx, value in enumerate(row_values):
            new_table.cell(r_idx, c_idx).text = value
    if style_id:
        _set_table_style_id(new_table, style_id)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return json.dumps({
        "status": "ok", "slide_index": slide_index, "shape_name": shape_name,
        "rows_written": n_rows, "cols_written": n_cols, "output_path": str(output_path),
    }, indent=2)


def _estimate_pptx_text_capacity(registry: ToolRegistry, args: dict[str, Any]) -> str:
    pptx = _require("pptx", "python-pptx>=1.0", "pptx")
    _require("PIL", "Pillow>=9.0", "pptx")
    from . import text_metrics

    path = registry._resolve_path(str(args["path"]))
    slide_index = int(args["slide_index"])
    shape_name = str(args["shape_name"])
    paragraphs = [str(p) for p in args.get("paragraphs", [])]
    is_chinese = args.get("is_chinese")
    if is_chinese is not None:
        is_chinese = bool(is_chinese)

    prs = pptx.Presentation(str(path))
    slide = _resolve_slide(prs, slide_index)
    shape = _find_shape_by_name(slide, shape_name)
    if not shape.has_text_frame:
        raise ValueError(
            f"Shape '{shape.name}' on slide {slide_index} has no text frame "
            f"(is_table={shape.has_table}) — capacity estimation only applies to text shapes."
        )

    try:
        result = text_metrics.estimate_capacity(shape, paragraphs, is_chinese=is_chinese)
    except FileNotFoundError as exc:
        raise ValueError(str(exc)) from exc

    result["slide_index"] = slide_index
    result["shape_name"] = shape.name
    return json.dumps(result, indent=2)


def _locate_databook_stage_columns(registry: ToolRegistry, args: dict[str, Any]) -> str:
    openpyxl = _require("openpyxl", "openpyxl>=3.1", "excel")
    from . import databook_metrics

    path = registry._resolve_path(str(args["path"]))
    sheet_name = args.get("sheet")
    max_scan_rows = int(args.get("max_scan_rows", 20))

    wb = openpyxl.load_workbook(str(path), data_only=True)
    ws = wb[sheet_name] if sheet_name else wb.active
    scan_end = min(ws.min_row + max_scan_rows, ws.max_row)
    rows = [[c.value for c in row_cells] for row_cells in ws.iter_rows(min_row=ws.min_row, max_row=scan_end)]
    resolved_sheet_name = ws.title
    wb.close()

    result = databook_metrics.locate_stage_columns(rows)
    result["sheet"] = resolved_sheet_name
    return json.dumps(result, indent=2, ensure_ascii=False, default=str)


# ---- Excel preview handler --------------------------------------------------

def _read_excel_preview(registry: ToolRegistry, args: dict[str, Any]) -> str:
    openpyxl = _require("openpyxl", "openpyxl>=3.1", "excel")
    path = registry._resolve_path(str(args["path"]))
    wb = openpyxl.load_workbook(str(path), data_only=True)
    sheet_name = args.get("sheet")
    ws = wb[sheet_name] if sheet_name else wb.active
    max_rows = int(args.get("max_rows", 20))
    hidden_rows, hidden_cols = _get_visibility_state(ws)
    rows: list[list[str]] = []
    for row_cells in ws.iter_rows():
        if not row_cells:
            continue
        if row_cells[0].row in hidden_rows:
            continue
        if len(rows) >= max_rows + 1:
            break
        rows.append([
            str(c.value) if c.value is not None else ""
            for c in row_cells if c.column not in hidden_cols
        ])
    wb.close()
    if not rows:
        return json.dumps({"headers": [], "rows": [], "schema": []})
    headers = rows[0]
    data = rows[1:]
    schema = []
    for col_idx, h in enumerate(headers):
        values = [r[col_idx] for r in data[:20] if col_idx < len(r) and r[col_idx]]
        numeric_count = sum(1 for v in values if _is_numeric(v))
        inferred = "numeric" if numeric_count > len(values) * 0.7 else "text"
        schema.append({"name": h, "inferred_type": inferred})
    return json.dumps({
        "headers": headers,
        "row_count": len(data),
        "truncated": len(data) >= max_rows,
        "schema": schema,
        "preview": data[:10],
    })


def _inspect_excel_sheets(registry: ToolRegistry, args: dict[str, Any]) -> str:
    openpyxl = _require("openpyxl", "openpyxl>=3.1", "excel")
    path = registry._resolve_path(str(args["path"]))
    wb = openpyxl.load_workbook(str(path), data_only=True)
    sample_rows = int(args.get("sample_rows", 3))
    sheets_info = []
    hidden_count = 0
    for ws in wb.worksheets:
        if ws.sheet_state != "visible":
            hidden_count += 1
            continue
        hidden_rows, hidden_cols = _get_visibility_state(ws)
        rows: list[list[str]] = []
        for row_cells in ws.iter_rows():
            if not row_cells:
                continue
            if row_cells[0].row in hidden_rows:
                continue
            rows.append([
                str(c.value) if c.value is not None else ""
                for c in row_cells if c.column not in hidden_cols
            ])
            if len(rows) >= sample_rows + 1:
                break
        headers = rows[0] if rows else []
        samples = rows[1:] if len(rows) > 1 else []
        total = max(0, (ws.max_row or 1) - 1 - len(hidden_rows))
        entry: dict[str, Any] = {
            "name": ws.title,
            "estimated_data_rows": total,
            "headers": headers,
            "sample": samples,
        }
        if hidden_rows:
            entry["hidden_rows"] = len(hidden_rows)
        if hidden_cols:
            entry["hidden_cols"] = len(hidden_cols)
        sheets_info.append(entry)
    wb.close()
    result: dict[str, Any] = {"sheets": sheets_info}
    if hidden_count:
        result["hidden_sheets"] = hidden_count
    return json.dumps(result, indent=2)


def _is_numeric(value: str) -> bool:
    try:
        float(value.replace(",", ""))
        return True
    except (ValueError, AttributeError):
        return False


# ---- Image handler ----------------------------------------------------------

def _image_read(registry: ToolRegistry, args: dict[str, Any]) -> str:
    path = registry._resolve_path(str(args["path"]))
    if not path.exists():
        raise FileNotFoundError(f"Image not found: {path}")
    suffix = path.suffix.lower()
    if suffix not in _IMAGE_SUFFIXES:
        raise ValueError(f"Unsupported image extension `{suffix}` (expected one of {list(_IMAGE_SUFFIXES)}).")
    size_bytes = path.stat().st_size
    mime, _ = mimetypes.guess_type(str(path))
    mime = mime or f"image/{suffix.lstrip('.')}"
    include_b64 = bool(args.get("include_base64", False))

    info: dict[str, Any] = {
        "path": str(path.relative_to(registry.workspace_root)) if registry.workspace_root in path.parents else str(path),
        "mime_type": mime,
        "size_bytes": size_bytes,
    }

    # Try to extract width/height/mode via PIL if available; degrade gracefully.
    try:
        from PIL import Image  # type: ignore
        with Image.open(str(path)) as im:
            info["width"] = im.width
            info["height"] = im.height
            info["mode"] = im.mode
    except ImportError:
        info["dimensions_note"] = "Install Pillow (`pip install pillow`) to extract width/height."
    except Exception as exc:
        info["dimensions_error"] = str(exc)

    if include_b64:
        if size_bytes > _IMAGE_MAX_BASE64_BYTES:
            info["base64_skipped"] = (
                f"Image is {size_bytes:,} bytes; exceeds {_IMAGE_MAX_BASE64_BYTES:,} byte cap. "
                "Resize the image or set include_base64=false."
            )
        else:
            raw = path.read_bytes()
            data_url = f"data:{mime};base64,{base64.b64encode(raw).decode('ascii')}"
            info["data_url"] = data_url

    return json.dumps(info, indent=2)


# ---- PDF handler ------------------------------------------------------------

def _read_pdf_text(registry: ToolRegistry, args: dict[str, Any]) -> str:
    pdfplumber = _require("pdfplumber", "pdfplumber>=0.10", "pdf")
    path = registry._resolve_path(str(args["path"]))
    max_pages = int(args.get("max_pages", 50))
    # layout=True approximates `pdftotext -layout`: whitespace preserves the
    # printed column positions, which financial-statement extraction needs
    # (and locked-down machines can't run the pdftotext executable).
    layout = bool(args.get("layout", False))
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages):
            if i >= max_pages:
                break
            text = page.extract_text(layout=layout) or ""
            pages.append({"page": i + 1, "text": text})
    return json.dumps({"pages": pages, "total_pages": len(pages)})


def _extract_pdf_words(registry: ToolRegistry, args: dict[str, Any]) -> str:
    _require("pdfplumber", "pdfplumber>=0.10", "pdf")
    from . import pdf_tables
    path = registry._resolve_path(str(args["path"]))
    words = pdf_tables.extract_page_words(str(path), int(args["page"]))
    return json.dumps({"page": int(args["page"]), "words": words, "count": len(words)})


def _extract_pdf_period_table(registry: ToolRegistry, args: dict[str, Any]) -> str:
    _require("pdfplumber", "pdfplumber>=0.10", "pdf")
    from . import pdf_tables
    path = registry._resolve_path(str(args["path"]))
    periods = args.get("periods")
    if periods is not None:
        periods = [str(p) for p in periods]
    result = pdf_tables.extract_period_table(str(path), int(args["page"]), periods)
    return json.dumps(result, ensure_ascii=False)
